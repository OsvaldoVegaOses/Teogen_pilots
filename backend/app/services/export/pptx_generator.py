from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any, Dict, Iterable, List

from pptx import Presentation


class PptxGenerator:
    @staticmethod
    def _as_text(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, (str, int, float, bool)):
            return str(value)
        if isinstance(value, list):
            return " | ".join([PptxGenerator._as_text(v) for v in value if v is not None])
        if isinstance(value, dict):
            preferred = [
                value.get("text"),
                value.get("description"),
                value.get("definition"),
                value.get("name"),
                value.get("id"),
            ]
            for item in preferred:
                if item:
                    return str(item)
            return str(value)
        return str(value)

    @staticmethod
    def _as_bullets(items: Iterable[Any], limit: int = 8) -> List[str]:
        if isinstance(items, (str, bytes, dict)):
            items = [items]
        lines: List[str] = []
        for item in items:
            text = PptxGenerator._as_text(item).strip()
            if text:
                lines.append(f"- {text}")
            if len(lines) >= limit:
                break
        return lines

    @staticmethod
    def _flatten_evidence(theory_data: Dict[str, Any], limit: int = 8) -> List[str]:
        validation = theory_data.get("validation", {}) or {}
        summary = validation.get("network_metrics_summary", {}) or {}
        semantic = summary.get("semantic_evidence_top", []) or []
        lines: List[str] = []

        for bucket in semantic:
            category = bucket.get("category_name", "Categoria")
            for frag in bucket.get("fragments", []) or []:
                score = frag.get("score")
                score_text = f" score={score:.3f}" if isinstance(score, (int, float)) else ""
                text = PptxGenerator._as_text(frag.get("text", "")).strip()
                if text:
                    lines.append(f"- [{category}]{score_text}: {text[:220]}")
                if len(lines) >= limit:
                    return lines
        return lines

    @staticmethod
    def _add_title_content_slide(prs: Presentation, title: str, lines: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "\n".join(lines) if lines else "Sin datos"

    def generate(self, project_name: str, theory_data: Dict[str, Any], template_key: str = "generic") -> BytesIO:
        prs = Presentation()

        model = theory_data.get("model_json", {}) or {}
        validation = theory_data.get("validation", {}) or {}
        summary = validation.get("network_metrics_summary", {}) or {}
        counts = summary.get("counts", {}) or {}
        propositions = theory_data.get("propositions", []) or []
        gaps = theory_data.get("gaps", []) or []

        central = self._as_text(
            model.get("selected_central_category")
            or (model.get("central_phenomenon", {}) or {}).get("name")
            or "No disponible"
        )

        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = f"TheoGen Report - {project_name}"
        cover.placeholders[1].text = (
            f"Template: {template_key}\n"
            f"Fecha: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        )

        self._add_title_content_slide(
            prs,
            "Resumen Ejecutivo",
            [
                f"- Categoria central: {central}",
                f"- Confianza: {theory_data.get('confidence_score', 0)}",
                f"- Proposiciones: {len(propositions)}",
                f"- Brechas: {len(gaps)}",
                f"- Entrevistas (grafo): {counts.get('fragment_count', 0)}",
                f"- Categorias (grafo): {counts.get('category_count', 0)}",
                f"- Generado por: {theory_data.get('generated_by', 'TheoGen')}",
            ],
        )

        self._add_title_content_slide(
            prs,
            "Categoria Central",
            [f"- {central}"]
            + self._as_bullets(model.get("central_category_justification", []), limit=5),
        )

        self._add_title_content_slide(
            prs,
            "Condiciones Causales",
            self._as_bullets(model.get("conditions", []), limit=8)
            or [f"- {self._as_text(model.get('conditions', 'No disponible'))}"],
        )

        self._add_title_content_slide(
            prs,
            "Contexto e Intervinientes",
            self._as_bullets(model.get("intervening_conditions", []), limit=8)
            or self._as_bullets(model.get("context", []), limit=8)
            or [f"- {self._as_text(model.get('context', 'No disponible'))}"],
        )

        self._add_title_content_slide(
            prs,
            "Acciones e Interacciones",
            self._as_bullets(model.get("actions", []), limit=8)
            or [f"- {self._as_text(model.get('actions', 'No disponible'))}"],
        )

        self._add_title_content_slide(
            prs,
            "Consecuencias",
            self._as_bullets(model.get("consequences", []), limit=8)
            or [f"- {self._as_text(model.get('consequences', 'No disponible'))}"],
        )

        self._add_title_content_slide(
            prs,
            "Proposiciones Teoricas",
            self._as_bullets(propositions, limit=12),
        )

        self._add_title_content_slide(
            prs,
            "Brechas y Riesgos",
            self._as_bullets(gaps, limit=12),
        )

        evidence_lines = self._flatten_evidence(theory_data, limit=10)
        self._add_title_content_slide(
            prs,
            "Evidencia Relevante",
            evidence_lines or ["- No se registraron fragmentos de evidencia en validacion."],
        )

        self._add_title_content_slide(
            prs,
            "Proximos Pasos",
            [
                "- Validar proposiciones con nuevos casos extremos.",
                "- Confirmar categorias perifericas con baja evidencia.",
                "- Priorizar acciones sobre brechas de mayor impacto.",
                "- Recalibrar el modelo tras nueva ingesta de entrevistas.",
            ],
        )

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf
